#!/usr/bin/env python3
"""
PowerPoint Processor

Extract text content from PPTX files including slide content,
speaker notes, and preserve slide structure.
"""

import asyncio
import logging
from pathlib import Path
from typing import Dict, Any, List
import time

from .base_processor import BaseFileProcessor, ProcessingResult

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PPTXProcessor(BaseFileProcessor):
    """PowerPoint PPTX file processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pptx']
        self.supported_mime_types = ['application/vnd.openxmlformats-officedocument.presentationml.presentation']
        self.pptx_available = PPTX_AVAILABLE
        
        if not self.pptx_available:
            self.logger.error("python-pptx library not available")
    
    async def extract_text(self, file_path: Path) -> ProcessingResult:
        """Extract text from PowerPoint presentation"""
        start_time = time.time()
        
        if not self.pptx_available:
            metadata = self._create_base_metadata(file_path, 0)
            metadata.errors.append("python-pptx library not available")
            return ProcessingResult(success=False, text_content="", metadata=metadata)
        
        try:
            prs = Presentation(file_path)
            slides_content = []
            structured_content = {'slides': []}
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_data = {'slide': slide_num, 'content': '', 'notes': '', 'shapes': 0}
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    slide_data['shapes'] += 1
                
                slide_content = "\n".join(slide_text)
                slide_data['content'] = slide_content
                
                # Extract speaker notes
                notes_text = ""
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    slide_data['notes'] = notes_text
                
                # Combine slide content
                full_slide_text = f"Slide {slide_num}:\n{slide_content}"
                if notes_text:
                    full_slide_text += f"\n\nNotes: {notes_text}"
                
                slides_content.append(full_slide_text)
                structured_content['slides'].append(slide_data)
            
            # Combine all text
            full_text = "\n\n".join(slides_content)
            normalized_text = self._normalize_text(full_text)
            
            # Create metadata
            processing_time_ms = int((time.time() - start_time) * 1000)
            metadata = self._create_base_metadata(file_path, processing_time_ms)
            metadata.page_count = len(prs.slides)
            metadata.character_count = len(normalized_text)
            metadata.word_count = self._count_words(normalized_text)
            metadata.language = self._detect_language(normalized_text)
            metadata.extraction_method = "python-pptx"
            
            return ProcessingResult(
                success=True,
                text_content=normalized_text,
                metadata=metadata,
                structured_content=structured_content
            )
            
        except Exception as e:
            processing_time_ms = int((time.time() - start_time) * 1000)
            metadata = self._create_base_metadata(file_path, processing_time_ms)
            metadata.errors.append(f"PPTX processing error: {str(e)}")
            
            return ProcessingResult(success=False, text_content="", metadata=metadata)